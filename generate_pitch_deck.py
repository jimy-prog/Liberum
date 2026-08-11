from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_pitch_deck():
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    title_content_layout = prs.slide_layouts[1]
    
    slides_data = [
        {
            "type": "title",
            "title": "Liberum",
            "subtitle": "The Future of Language Education Management\n\nPitch Deck"
        },
        {
            "type": "content",
            "title": "Slide 2: The Origin Story",
            "content": "• Built out of pure frustration with existing tools.\n• Bouncing between WhatsApp, Excel, and paper mock exams.\n• Realized millions of educators face the same fragmentation.\n• Built Liberum to be the ultimate unified ecosystem."
        },
        {
            "type": "content",
            "title": "Slide 3: The Problem (The 'Frankenstein' Setup)",
            "content": "• Language centers patch together 5+ generic tools to survive.\n• Leads to lost revenue and hundreds of hours wasted on manual admin work.\n• Disconnected and frustrating experience for students.\n• The market desperately needs a single, unified solution."
        },
        {
            "type": "content",
            "title": "Slide 4: The Solution (Liberum MVP)",
            "content": "• A unified dashboard for center owners to manage their entire business.\n• Group management, automated billing, and timetable scheduling.\n• Simple, powerful, and it works today."
        },
        {
            "type": "content",
            "title": "Slide 5: The Unfair Advantage (AI Mock Exams)",
            "content": "• Proprietary AI-powered Mock Exam engine built directly into the platform.\n• Students take full IELTS exams, and AI grades them instantly.\n• Covers reading, listening, writing, and speaking.\n• Saves thousands of dollars in teacher grading hours."
        },
        {
            "type": "content",
            "title": "Slide 6: The Grand Vision (The Ecosystem)",
            "content": "• We are building a complete educational ecosystem:\n  - Liberum Studio & Shorts: Interactive, bite-sized content creation.\n  - Liberum Mock: A standalone platform for global AI test prep.\n  - Liberum Market: A global marketplace for tutors to sell courses."
        },
        {
            "type": "content",
            "title": "Slide 7: Business Model & Pricing",
            "content": "• B2B SaaS subscription based on student volume (e.g., $X/month per 100 students).\n• Premium add-ons for AI mock exam credits.\n• Highly scalable recurring revenue stream."
        },
        {
            "type": "content",
            "title": "Slide 8: Profitability & Market Potential",
            "content": "• Extremely low infrastructure costs = incredibly high profit margins.\n• Projected to reach profitability within [X] months of launch.\n• High Lifetime Value (LTV) per customer as we roll out Studio & Market."
        },
        {
            "type": "content",
            "title": "Slide 9: Current Traction",
            "content": "• Actively running a successful beta test with a real language center.\n• Managing real students and processing real AI mock tests.\n• Phenomenal feedback proving actual product-market fit."
        },
        {
            "type": "content",
            "title": "Slide 10: Why Invest Now?",
            "content": "• Language education technology is stuck in the past; the market is ready.\n• We have the right product, vision, and traction to dominate.\n• Raising [Amount] for GTM strategy, AI R&D, and expanding the ecosystem.\n• Join us in building the future of education."
        }
    ]
    
    for slide_data in slides_data:
        if slide_data["type"] == "title":
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = slide_data["title"]
            subtitle.text = slide_data["subtitle"]
        else:
            slide = prs.slides.add_slide(title_content_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = slide_data["title"]
            content.text = slide_data["content"]
            
    # Save the presentation to the user's desktop
    output_path = "/Users/jamshidmahkamov/Desktop/Liberum_Investor_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_pitch_deck()
